#this will use a library to detect the language of the document
#it will return the language of the document

from langdetect import detect
from pathlib import Path

def extract_text(file_path: str) -> str:
    """Extract text from PDF, DOCX, PPTX, or TXT."""
    file_ext = Path(file_path).suffix.lower()
    text = ""

    if file_ext == ".pdf":
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

    elif file_ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif file_ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    else:
        # fallback: treat as plain text
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()

    return text.strip()


def get_file_language(file_path: str) -> str:
    """Detect the language of the document."""
    text = extract_text(file_path)
    if not text:
        return "unknown"
    try:
        language = detect(text)
        return language, text
    except Exception:
        return "unknown"
